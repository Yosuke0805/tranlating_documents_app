''' coding: utf-8 '''
# ------------------------------------------------------------
# Content : Creating a tool for automcatically tranlating documents
# Author : Yosuke Kawazoe
# Data Updated：
# Update Details：
# ------------------------------------------------------------

# Import
import os
import streamlit as st
import traceback
import requests
import pandas as pd
import openpyxl
import tempfile
from pptx import Presentation
from googletrans import Translator
from deep_translator import GoogleTranslator


# ------------------------------------------------------------
# ★★★★★★  main part ★★★★★★
# ------------------------------------------------------------
def main():

    try:
        # Set the title and description
        st.title("Translating you file into English ")


        # File uploader widget
        uploaded_file = st.file_uploader("Upload a PowerPoint (.pptx) or Excel (.xlsx) file", type=["pptx", "xlsx"])
        dest_lang = st.text_input("Enter the destination language (e.g., 'en' for English, 'ja' for Japanese):", "en")

        if st.button('Translate'):
            if uploaded_file is not None:
                with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                    temp_file.write(uploaded_file.read())
                    temp_file_path = temp_file.name

                if uploaded_file.name.endswith(".pptx"):
                    translated_file_path = translate_presentation(temp_file_path, dest_lang)
                elif uploaded_file.name.endswith(".xlsx"):
                    translated_file_path = translate_excel(temp_file_path, dest_lang)

                with open(translated_file_path, "rb") as file:
                    btn = st.download_button(
                        label="Download Translated File",
                        data=file,
                        file_name=f"translated_{uploaded_file.name}",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation" if uploaded_file.name.endswith(".pptx") else "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    )

                os.remove(temp_file_path)
                os.remove(translated_file_path)
            else:
                st.error("Please upload a file to translate.")
    except Exception:
        traceback.print_exc()


# ------------------------------------------------------------
# ★★★★★★  functions ★★★★★★
# ------------------------------------------------------------
def translate_text(text, translator):
    if text.strip() == "":
        return text
    try:
        translation = translator.translate(text, src='ja', dest='en')
        return translation.text
    except Exception as e:
        print(f"Error translating text '{text}': {e}")
        return text

def translate_presentation(input_path, output_path):
    # Load the presentation
    prs = Presentation(input_path)
    translator = Translator()

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        translated_text = translate_text(original_text, translator)
                        run.text = translated_text
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text
                                translated_text = translate_text(original_text, translator)
                                run.text = translated_text

    # Save the translated presentation
    prs.save(output_path)


def translate_excel(input_path, output_path):
    # Load the Excel workbook
    workbook = openpyxl.load_workbook(input_path)
    translator = GoogleTranslator(source='ja', target='en')

    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None and isinstance(cell.value, str):
                    original_text = cell.value
                    translated_text = translate_text(original_text, translator)
                    cell.value = translated_text

    # Save the translated workbook
    workbook.save(output_path)

# ------------------------------------------------------------
# ★★★★★★  execution part  ★★★★★★
# ------------------------------------------------------------
if __name__ == '__main__':

    # execute
    main()
